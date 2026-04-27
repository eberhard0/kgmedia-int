"""
KG Media Amplification Watch — editorial team briefing deck.

Targeted at non-technical editorial readers. No jargon. Each slide
either explains a concept in plain language or walks through one part
of the dashboard.

Output:
  /home/eberhard/projects/kgmedia-int/docs/amplification-editorial-briefing.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Color palette (mirrors the dashboard)
BG          = RGBColor(0x0B, 0x0F, 0x19)
CARD        = RGBColor(0x0F, 0x17, 0x2A)
CARD_BORDER = RGBColor(0x1F, 0x29, 0x37)
TEXT        = RGBColor(0xF8, 0xFA, 0xFC)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
DIM         = RGBColor(0x64, 0x74, 0x8B)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
RED         = RGBColor(0xEF, 0x44, 0x44)
RED_TINT    = RGBColor(0x7F, 0x1D, 0x1D)
YELLOW      = RGBColor(0xFA, 0xCC, 0x15)
YELLOW_TINT = RGBColor(0x71, 0x4A, 0x10)
GREEN       = RGBColor(0x4A, 0xDE, 0x80)

# Platform tile colors (from app/amplification/page.tsx PLATFORM_META)
PLATFORM_COLORS = {
    "Google News": (RGBColor(0x60, 0xA5, 0xFA), "G"),
    "Reddit":      (RGBColor(0xFB, 0x92, 0x3C), "R"),
    "X / Twitter": (RGBColor(0x38, 0xBD, 0xF8), "X"),
    "TikTok":      (RGBColor(0xF4, 0x72, 0xB6), "♪"),
    "Instagram":   (RGBColor(0xC0, 0x84, 0xFC), "◉"),
    "Facebook":    (RGBColor(0x81, 0x8C, 0xF8), "f"),
}

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG


def add_text_box(slide, x, y, w, h, text, *,
                 size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Helvetica Neue"
    return tb


def add_card(slide, x, y, w, h, fill=CARD, border=CARD_BORDER, line_w=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.06
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = line_w
    return s


def add_footer(slide, page_num, total_pages):
    add_text_box(
        slide,
        Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        "KG Media Amplification Watch · editorial briefing",
        size=9, color=DIM,
    )
    add_text_box(
        slide,
        Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
        f"{page_num} / {total_pages}",
        size=9, color=DIM, align=PP_ALIGN.RIGHT,
    )


# ────────────────── SLIDES ──────────────────

PAGES = []  # populated; len() used for footer numbering


def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.5),
        "KG Media Amplification Watch",
        size=54, bold=True, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.0),
        "What it is, how to read it, and what to do when it fires",
        size=22, color=ACCENT_BLUE,
    )
    add_text_box(
        s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.6),
        "For the editorial team   ·   April 2026",
        size=14, color=MUTED,
    )
    PAGES.append(s)


def slide_one_question():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.7),
        "The one question this dashboard answers",
        size=28, bold=True, color=TEXT,
    )
    add_card(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.0),
             fill=RGBColor(0x16, 0x1B, 0x2E))
    add_text_box(
        s, Inches(2.0), Inches(2.5), Inches(9.3), Inches(1.5),
        '"Is anyone outside KG Media talking about us right now,'
        ' and how loud is it?"',
        size=24, color=ACCENT_BLUE, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.5),
        [
            "Every post you'll see on this page is from someone else — a TikTok",
            "user, a Reddit thread, a news outlet, a Facebook page — and it",
            "explicitly mentions a KG Media brand (Kompas, Kompas.com, Kompas.id,",
            "Harian Kompas, KompasTV, Kompasiana, Kontan, Gramedia, Santika).",
            "",
            "If a post doesn't name one of those brands, it doesn't appear here.",
        ],
        size=16, color=TEXT,
    )
    PAGES.append(s)


def slide_why_it_exists():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.7),
        "Why this exists",
        size=28, bold=True, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.0),
        "Stories about KG brands can travel from a single Tweet to a national",
        size=18, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(2.05), Inches(11.7), Inches(1.0),
        "conversation in hours. Editorial used to find out from phone calls,",
        size=18, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.0),
        "screenshots, or after the wave had already broken.",
        size=18, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.6),
        "This page is the early-warning radar:",
        size=18, color=TEXT, bold=True,
    )
    bullets = [
        "•   Watches 6 platforms automatically every hour",
        "•   Counts mentions of any KG brand",
        "•   Groups them into 'stories' so you don't read the same thing 50 times",
        "•   Flashes red when a story crosses the threshold for editorial attention",
    ]
    for i, b in enumerate(bullets):
        add_text_box(
            s, Inches(1.2), Inches(4.3 + i * 0.55), Inches(11.0), Inches(0.5),
            b, size=16, color=TEXT,
        )
    PAGES.append(s)


def slide_60s_routine():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "Your 60-second daily routine",
        size=28, bold=True, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
        "Open the page once a day. Read top to bottom. Stop as soon as nothing's wrong.",
        size=14, color=MUTED,
    )
    steps = [
        ("1", "Red banner at top?",
         "Critical alert. Click into it. Read 3-5 mentions. Decide: clarify, follow up, or note it.",
         RED),
        ("2", "Top stories now panel?",
         "Three biggest stories ranked. Click any row to jump to its detail card below.",
         ACCENT_BLUE),
        ("3", "Trending section?",
         "Stories climbing toward critical. Skim entity names; investigate brand-aligned ones.",
         YELLOW),
        ("4", "Nothing? Close the tab.",
         "If the alert banner is empty and trending is quiet, there's nothing to act on. Done.",
         GREEN),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        y = Inches(1.95 + i * 1.25)
        # circle with step number
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y, Inches(0.7), Inches(0.7))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        ctf = circ.text_frame; ctf.margin_left = ctf.margin_right = Inches(0)
        ctf.margin_top = ctf.margin_bottom = Inches(0)
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num; cr.font.size = Pt(22); cr.font.bold = True
        cr.font.color.rgb = BG
        # title + body
        add_text_box(s, Inches(1.85), y, Inches(11), Inches(0.5),
                     title, size=18, bold=True, color=color)
        add_text_box(s, Inches(1.85), y + Inches(0.45), Inches(11), Inches(0.6),
                     body, size=14, color=TEXT)
    PAGES.append(s)


def slide_two_tiers():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "Two tiers — Trending and Critical",
        size=28, bold=True, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
        "Both tiers require at least 3 different sources. Mention count separates them.",
        size=14, color=MUTED,
    )

    # Trending card
    add_card(s, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.5),
             fill=YELLOW_TINT, border=YELLOW)
    add_text_box(s, Inches(1.1), Inches(2.3), Inches(5.2), Inches(0.6),
                 "TRENDING", size=22, bold=True, color=YELLOW)
    add_text_box(s, Inches(1.1), Inches(2.95), Inches(5.2), Inches(0.5),
                 "10–99 mentions in 24 hours", size=15, color=TEXT)
    add_text_box(s, Inches(1.1), Inches(3.6), Inches(5.2), Inches(2.5),
                 [
                     "Watch list — story is climbing,",
                     "not yet a public-facing alert.",
                     "",
                     "No pulse, no banner.",
                     "",
                     "Worth a look.",
                     "Not urgent.",
                 ], size=14, color=TEXT)

    # Critical card
    add_card(s, Inches(6.9), Inches(2.1), Inches(5.8), Inches(4.5),
             fill=RED_TINT, border=RED)
    add_text_box(s, Inches(7.2), Inches(2.3), Inches(5.2), Inches(0.6),
                 "CRITICAL", size=22, bold=True, color=RED)
    add_text_box(s, Inches(7.2), Inches(2.95), Inches(5.2), Inches(0.5),
                 "100+ mentions in 24 hours", size=15, color=TEXT)
    add_text_box(s, Inches(7.2), Inches(3.6), Inches(5.2), Inches(2.5),
                 [
                     "Real public conversation",
                     "around a KG brand.",
                     "",
                     "Red border + pulsing banner",
                     "on the dashboard.",
                     "",
                     "Editorial attention warranted.",
                 ], size=14, color=TEXT)
    add_text_box(
        s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
        "Why \"3 sources\" matters: 100 reposts of one viral TikTok still counts as one source. We want genuine cross-platform conversation.",
        size=12, color=MUTED,
    )
    PAGES.append(s)


def slide_what_a_cluster_is():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "What is a cluster?",
        size=28, bold=True, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.7),
        "A cluster = a group of mentions across platforms that all appear to be about the same story.",
        size=16, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.5),
        "Example:",
        size=16, bold=True, color=ACCENT_BLUE,
    )

    # Visual: source story → cluster of mentions
    add_card(s, Inches(0.8), Inches(2.6), Inches(4.0), Inches(2.0),
             fill=RGBColor(0x16, 0x1B, 0x2E))
    add_text_box(s, Inches(1.0), Inches(2.75), Inches(3.6), Inches(0.4),
                 "KOMPAS HEADLINE", size=10, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(1.0), Inches(3.15), Inches(3.6), Inches(1.4),
                 ['"Jemaah haji Indonesia',
                  ' meninggal di Madinah"'],
                 size=14, color=TEXT)

    # arrow
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.0), Inches(3.4), Inches(0.7), Inches(0.4))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED; arrow.line.fill.background()

    # mentions box
    add_card(s, Inches(5.9), Inches(2.6), Inches(6.8), Inches(4.0),
             fill=RGBColor(0x16, 0x1B, 0x2E), border=RED)
    add_text_box(s, Inches(6.1), Inches(2.75), Inches(6.4), Inches(0.4),
                 "CLUSTER (28 mentions, 5 sources)", size=10, bold=True, color=RED)
    sample_mentions = [
        "Reddit (r/indonesia): \"Sad news from Madinah...\"",
        "TikTok (@hariankompas re-shared): \"Innalillahi...\"",
        "Google News (Kompas.com): \"Satu Jemaah Haji...\"",
        "Facebook (Kompas page): comments piling up",
        "Reddit (r/Islam): linked discussion",
        "...22 more mentions",
    ]
    for i, m in enumerate(sample_mentions):
        add_text_box(s, Inches(6.1), Inches(3.2 + i * 0.45), Inches(6.4), Inches(0.4),
                     "• " + m, size=12, color=TEXT)

    add_text_box(
        s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
        "All those mentions get grouped into one card. You read one card, not 28 things.",
        size=13, color=MUTED,
    )
    PAGES.append(s)


def slide_dashboard_anatomy():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "Anatomy of the dashboard",
        size=28, bold=True, color=TEXT,
    )

    # Mockup of the page layout, top to bottom
    layout_items = [
        ("Page header  ·  KG Media Amplification Watch", 0.45, ACCENT_BLUE, BG),
        ("🚨  Critical alert banner  (only when there's a critical cluster)", 0.55, RED, RED_TINT),
        ("Top stories now  ·  the 3 biggest live clusters", 0.55, ACCENT_BLUE, RGBColor(0x14,0x1F,0x35)),
        ("Platform tiles  ·  Google News  ·  Reddit  ·  X  ·  TikTok  ·  Instagram  ·  Facebook", 0.7, TEXT, CARD),
        ("Critical clusters  (red, pulsing border)", 0.7, RED, RGBColor(0x21,0x10,0x10)),
        ("Trending clusters  (yellow, no pulse)", 0.7, YELLOW, RGBColor(0x1F,0x1A,0x06)),
        ("Latest mentions  ·  raw feed at the bottom", 0.5, MUTED, CARD),
    ]
    y = Inches(1.5)
    for label, h_in, color, fill in layout_items:
        add_card(s, Inches(0.8), y, Inches(11.7), Inches(h_in), fill=fill)
        add_text_box(s, Inches(1.1), y, Inches(11.0), Inches(h_in),
                     label, size=14, bold=True, color=color,
                     anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(h_in + 0.12)
    PAGES.append(s)


def slide_top_stories():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "Top stories now — your shortcut",
        size=28, bold=True, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
        "If you only have 30 seconds, this is what to read.",
        size=14, color=MUTED,
    )

    # blue-bordered panel mockup
    add_card(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.6),
             fill=RGBColor(0x10, 0x1A, 0x33), border=ACCENT_BLUE)
    add_text_box(s, Inches(1.1), Inches(2.15), Inches(8), Inches(0.4),
                 "TOP STORIES NOW · 3", size=12, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(1.1), Inches(2.55), Inches(8), Inches(0.4),
                 "what to look at first", size=10, color=DIM)

    rows = [
        ("1.", "CRITICAL", RED, "Madinah",
         "142 mentions · 5 sources · Satu Jemaah Haji Indonesia Meninggal..."),
        ("2.", "TRENDING", YELLOW, "Bandara",
         "52 mentions · 4 sources · Virus Nipah Jadi Alarm Global..."),
        ("3.", "TRENDING", YELLOW, "ASEAN",
         "28 mentions · 7 sources · ASEAN Hopes to Transform..."),
    ]
    for i, (num, tier, tcolor, entity, meta) in enumerate(rows):
        y = Inches(3.2 + i * 0.85)
        add_card(s, Inches(1.1), y, Inches(11.1), Inches(0.65),
                 fill=RGBColor(0x14, 0x1F, 0x36), border=CARD_BORDER)
        add_text_box(s, Inches(1.25), y + Inches(0.05), Inches(0.4), Inches(0.55),
                     num, size=14, color=DIM, bold=True)
        # tier badge
        add_card(s, Inches(1.7), y + Inches(0.13), Inches(1.2), Inches(0.4),
                 fill=tcolor, border=tcolor)
        add_text_box(s, Inches(1.7), y + Inches(0.15), Inches(1.2), Inches(0.36),
                     tier, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
        # entity
        add_text_box(s, Inches(3.0), y + Inches(0.15), Inches(2.5), Inches(0.4),
                     entity, size=14, bold=True, color=TEXT)
        # meta
        add_text_box(s, Inches(5.5), y + Inches(0.18), Inches(6.5), Inches(0.4),
                     meta, size=11, color=MUTED)

    add_text_box(
        s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.5),
        "Tap any row → page jumps to the full cluster card and opens it. Empty panel = nothing to act on right now.",
        size=12, color=MUTED,
    )
    PAGES.append(s)


def slide_when_to_act():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "When to act, when to ignore",
        size=28, bold=True, color=TEXT,
    )
    rows = [
        ("Red banner pulsing at top",          "Always investigate, even if it might be a false positive.",          RED),
        ("Trending cluster, brand-aligned entity (a Kompas reporter, recent headline, KG figure)", "Investigate. Catch it before it becomes critical.", YELLOW),
        ("Trending cluster, generic entity (\"Indonesia\", \"data\", \"people\")", "Probably noise. Skip.", MUTED),
        ("Single platform tile is high but unrelated to a cluster", "Curiosity, not urgent. Drill in if you have time.", MUTED),
        ("Page is empty / quiet", "Nothing on KG is amplifying. Close the tab.", GREEN),
    ]
    y = Inches(1.6)
    for trigger, action, color in rows:
        add_card(s, Inches(0.8), y, Inches(11.7), Inches(0.95),
                 fill=CARD, border=CARD_BORDER)
        add_card(s, Inches(0.8), y, Inches(0.18), Inches(0.95),
                 fill=color, border=color)
        add_text_box(s, Inches(1.15), y + Inches(0.1), Inches(5.6), Inches(0.4),
                     trigger, size=13, bold=True, color=TEXT)
        add_text_box(s, Inches(1.15), y + Inches(0.5), Inches(11.0), Inches(0.4),
                     action, size=12, color=MUTED)
        y += Inches(1.05)
    PAGES.append(s)


def slide_data_source():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "Where the data comes from",
        size=28, bold=True, color=TEXT,
    )
    add_text_box(
        s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
        "We watch 6 sources every hour. Each tile shows how many mentions came from that source in the last 24 hours.",
        size=14, color=MUTED,
    )

    # platform tiles in a 3x2 grid
    items = list(PLATFORM_COLORS.items())
    cols = 3
    tile_w = Inches(3.85); tile_h = Inches(1.95)
    x0 = Inches(0.8); y0 = Inches(2.1); gap_x = Inches(0.15); gap_y = Inches(0.2)
    sample_counts = [221, 200, 15, 153, 8, 133]
    for i, (label, (color, glyph)) in enumerate(items):
        r = i // cols; c = i % cols
        x = x0 + (tile_w + gap_x) * c
        y = y0 + (tile_h + gap_y) * r
        add_card(s, x, y, tile_w, tile_h, border=color, fill=CARD)
        # glyph
        gly = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x + Inches(0.2), y + Inches(0.2), Inches(0.45), Inches(0.45))
        gly.fill.solid(); gly.fill.fore_color.rgb = color; gly.line.fill.background()
        gly.adjustments[0] = 0.18
        gtf = gly.text_frame; gtf.margin_left = gtf.margin_right = Inches(0)
        gtf.margin_top = gtf.margin_bottom = Inches(0); gtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        gp = gtf.paragraphs[0]; gp.alignment = PP_ALIGN.CENTER
        gr = gp.add_run(); gr.text = glyph; gr.font.size = Pt(20); gr.font.bold = True
        gr.font.color.rgb = BG
        # label
        add_text_box(s, x + Inches(0.75), y + Inches(0.2), tile_w - Inches(0.85), Inches(0.45),
                     label, size=14, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        # count (sample)
        add_text_box(s, x + Inches(0.2), y + Inches(0.85), tile_w - Inches(0.4), Inches(0.5),
                     f"{sample_counts[i]} mentions",
                     size=18, bold=True, color=TEXT)
        # caption
        add_text_box(s, x + Inches(0.2), y + Inches(1.35), tile_w - Inches(0.4), Inches(0.5),
                     "(tap to drill in)", size=10, color=DIM)
    PAGES.append(s)


def slide_what_to_do_when():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "What to do when a cluster fires",
        size=28, bold=True, color=TEXT,
    )
    steps = [
        ("Read the linked Kompas article first",
         "The cluster card shows the original Kompas story it's tied to. Click it. Skim it. That's the substance."),
        ("Sample 3-5 of the external mentions",
         "Click '+ details' to expand. Scroll the snippets. You'll see the angle each platform is taking — supportive, critical, mocking, factual."),
        ("Decide the action",
         "Three options: a follow-up piece, a clarification / response, or a desk note for editorial. Most clusters resolve themselves; not every cluster needs a piece."),
        ("False positive? Move on.",
         "Sometimes the cluster is just three random posts that happen to share a word. The auto-clustering isn't perfect — your judgment overrides it."),
    ]
    for i, (title, body) in enumerate(steps):
        y = Inches(1.7 + i * 1.3)
        add_card(s, Inches(0.8), y, Inches(11.7), Inches(1.15),
                 fill=CARD, border=CARD_BORDER)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(1.0), y + Inches(0.3), Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT_BLUE
        circ.line.fill.background()
        ctf = circ.text_frame; ctf.margin_left = ctf.margin_right = Inches(0)
        ctf.margin_top = ctf.margin_bottom = Inches(0)
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(i + 1); cr.font.size = Pt(18); cr.font.bold = True
        cr.font.color.rgb = BG
        # text
        add_text_box(s, Inches(1.8), y + Inches(0.15), Inches(10.5), Inches(0.4),
                     title, size=15, bold=True, color=TEXT)
        add_text_box(s, Inches(1.8), y + Inches(0.55), Inches(10.5), Inches(0.55),
                     body, size=12, color=MUTED)
    PAGES.append(s)


def slide_what_it_isnt():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "What this is NOT",
        size=28, bold=True, color=TEXT,
    )
    items = [
        ("Not sentiment analysis on the brand",
         "We don't (yet) tell you whether the chatter is positive or negative. Open mentions to read the tone."),
        ("Not a complete view of social",
         "We watch 6 platforms with public data. Telegram, WhatsApp, private groups — not visible here."),
        ("Not real-time",
         "Refreshes hourly. A story breaking at 2pm WIB shows up around 3pm WIB at the earliest."),
        ("Not a replacement for editorial judgment",
         "It surfaces signal. You decide what's a story."),
    ]
    for i, (title, body) in enumerate(items):
        y = Inches(1.7 + i * 1.3)
        add_card(s, Inches(0.8), y, Inches(11.7), Inches(1.15),
                 fill=CARD, border=CARD_BORDER)
        # red X mark area
        add_card(s, Inches(0.8), y, Inches(0.18), Inches(1.15), fill=RED, border=RED)
        add_text_box(s, Inches(1.15), y + Inches(0.15), Inches(11.0), Inches(0.4),
                     title, size=15, bold=True, color=TEXT)
        add_text_box(s, Inches(1.15), y + Inches(0.55), Inches(11.0), Inches(0.55),
                     body, size=12, color=MUTED)
    PAGES.append(s)


def slide_glossary():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7),
        "Quick glossary",
        size=28, bold=True, color=TEXT,
    )
    pairs = [
        ("Mention", "One post (a tweet, TikTok video, Reddit thread, news article) that names a KG brand."),
        ("Source", "The platform / outlet a mention came from. Reddit and TikTok are different sources; two TikTok videos are the same source."),
        ("Cluster", "A group of mentions about the same story, auto-detected. The cluster card is what you actually read."),
        ("Entity", "The main person, place, or organization the cluster is about. Auto-extracted from the source article."),
        ("Tier", "Trending = 10+ mentions / 3+ sources. Critical = 100+ mentions / 3+ sources. Critical pulses red."),
        ("24h window", "Everything on the page is from the last 24 hours. Older mentions roll off automatically."),
    ]
    y0 = Inches(1.6)
    for i, (term, defn) in enumerate(pairs):
        y = y0 + Inches(i * 0.85)
        add_text_box(s, Inches(0.8), y, Inches(2.5), Inches(0.4),
                     term, size=15, bold=True, color=ACCENT_BLUE)
        add_text_box(s, Inches(3.4), y, Inches(9.5), Inches(0.8),
                     defn, size=13, color=TEXT)
    PAGES.append(s)


def slide_questions():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text_box(
        s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5),
        "Questions?",
        size=72, bold=True, color=TEXT, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.6),
        "Open the FAQ on the dashboard for the long version.",
        size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.6),
        "https://kgmedia-int.vercel.app/faq#amplification",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )
    PAGES.append(s)


# Build slides in order
slide_title()
slide_one_question()
slide_why_it_exists()
slide_60s_routine()
slide_two_tiers()
slide_what_a_cluster_is()
slide_dashboard_anatomy()
slide_top_stories()
slide_when_to_act()
slide_data_source()
slide_what_to_do_when()
slide_what_it_isnt()
slide_glossary()
slide_questions()

# Now add footers (skip title and final 'Questions' slide for cleaner look)
total = len(PAGES)
for i, slide in enumerate(PAGES, start=1):
    if i in (1, total):
        continue
    add_footer(slide, i, total)

OUT_DIR = "/home/eberhard/projects/kgmedia-int/docs"
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, "amplification-editorial-briefing.pptx")
prs.save(out_path)
print(f"Wrote {out_path}  ({total} slides)")
